from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_box = slide.shapes.title
        title_box.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(44)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 136)
        
        if subtitle:
            subtitle_box = slide.placeholders[1]
            subtitle_box.text = subtitle
            subtitle_box.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(160, 160, 160)
        return slide
    
    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.title
        title_box.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(36)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 255, 136)
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(12.333)
        
        for i, item in enumerate(content_items):
            text_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
            text_frame = text_box.text_frame
            text_frame.text = item
            text_frame.paragraphs[0].font.size = Pt(18)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            text_frame.word_wrap = True
            top += Inches(0.5)
        
        return slide
    
    slide1 = add_title_slide(
        "足球教练助手",
        "AI驱动的智能训练分析系统"
    )
    
    slide2 = add_content_slide(
        "系统概述",
        [
            "智能化: 基于AI的姿势识别和动作分析",
            "数据化: 将训练过程转化为可量化的数据",
            "个性化: 根据用户数据定制训练计划",
            "游戏化: 成就系统激励持续训练",
            "目标用户: 足球运动员、足球教练、体育训练机构",
            "解决的问题: 训练动作不规范、数据难以追踪、缺乏个性化计划"
        ]
    )
    
    slide3 = add_content_slide(
        "核心功能 - 视频分析",
        [
            "姿势识别: 实时检测运动员的姿势状态",
            "动作分析: 分析运动轨迹和速度",
            "准确度评估: 量化动作的准确性(0-100%)",
            "问题诊断: 自动识别常见问题并提供改进建议",
            "技术实现: React + TypeScript / Python + FastAPI / MediaPipe",
            "用户价值: 节省时间、提高准确度、持续改进"
        ]
    )
    
    slide4 = add_content_slide(
        "核心功能 - 训练计划",
        [
            "训练日历: 可视化展示训练计划",
            "智能推荐: 基于用户能力水平推荐训练",
            "进度追踪: 记录每次训练的完成情况",
            "目标管理: 设置每周训练目标和准确度目标",
            "训练类型: 速度训练、力量训练、敏捷训练、技术训练",
            "用户价值: 结构化、目标导向、可视化"
        ]
    )
    
    slide5 = add_content_slide(
        "核心功能 - 数据分析",
        [
            "训练历史: 完整的训练记录，支持筛选和导出",
            "统计分析: 月度训练趋势、准确度变化曲线",
            "技能评估: 速度、平衡、协调、力量四大维度",
            "成就系统: 自动检测成就，激励持续训练",
            "数据可视化: 趋势图、柱状图、进度条、成就徽章",
            "用户价值: 深度洞察、趋势分析、目标导向"
        ]
    )
    
    slide6 = add_content_slide(
        "核心功能 - 用户档案",
        [
            "个人信息: 姓名、邮箱、电话、所在地、加入时间",
            "成就展示: 获得的成就列表和成就时间",
            "训练统计: 总训练次数、平均准确度、训练时长",
            "等级系统: 新手、进阶、专家、大师",
            "等级体系: 0-10次新手、11-50次进阶、51-100次专家、100+次大师",
            "用户价值: 个性化、成就感、数据驱动"
        ]
    )
    
    slide7 = add_content_slide(
        "核心功能 - 设置系统",
        [
            "外观设置: 主题切换、语言选择",
            "通知设置: 训练提醒、成就通知、每周报告",
            "训练目标: 每周训练目标设置、准确度目标设置",
            "隐私设置: 数据共享控制、使用分析开关",
            "数据管理: 导出设置、导入设置、重置为默认、清除所有数据",
            "用户价值: 个性化、数据安全、隐私保护"
        ]
    )
    
    slide8 = add_content_slide(
        "技术架构",
        [
            "前端技术栈: React 18 + TypeScript + Tailwind CSS + Vite",
            "后端技术栈: FastAPI + Python 3.x + MediaPipe + WebSocket",
            "数据管理: DataManager单例、SettingsManager单例、LocalStorage",
            "整体架构: 前后端分离、HTTP/WebSocket通信、LocalStorage存储",
            "核心优势: 高性能、类型安全、现代化UI、AI驱动、实时通信"
        ]
    )
    
    slide9 = add_content_slide(
        "用户体验设计",
        [
            "设计理念: 足球运动美学、黑色为主、绿色为辅",
            "配色方案: 深黑背景(#0a0a0a)、霓虹绿强调(#00ff88)",
            "动画效果: 背景渐变、脉冲光晕、交互动画、装饰动画",
            "玻璃态设计: 半透明卡片、20px背景模糊、半透明白色边框",
            "响应式设计: 桌面、平板、移动端适配",
            "用户价值: 视觉享受、沉浸体验、多端适配、流畅交互"
        ]
    )
    
    slide10 = add_content_slide(
        "欢迎页面",
        [
            "设计特点: 世界杯主题、奖杯和足球元素、进球集锦",
            "动画系统: 多层脉冲环、背景动画、装饰性脉冲点",
            "自动播放: 每5秒自动切换经典进球、无需用户交互",
            "进入按钮: 2秒后显示、浮动动画、绿色渐变",
            "进球集锦: 梅西(2022)、C罗(2018)、内马尔(2014)、姆巴佩(2018)、贝克汉姆(2001)",
            "用户价值: 视觉冲击、足球氛围、沉浸体验、快速进入"
        ]
    )
    
    slide11 = add_content_slide(
        "技术亮点",
        [
            "AI驱动的姿势识别: MediaPipe Pose Detection、实时检测33个身体关键点",
            "实时数据处理: WebSocket通信、低延迟数据传输、实时更新前端",
            "单例模式数据管理: 全局唯一数据实例、自动同步所有组件",
            "暗黑主题系统: CSS变量、主题切换无刷新、玻璃态设计",
            "响应式设计: Tailwind CSS、移动优先设计、优秀的移动端体验",
            "类型安全: TypeScript编译时类型检查、智能代码提示、减少运行时错误"
        ]
    )
    
    slide12 = add_content_slide(
        "未来规划",
        [
            "短期目标(3-6个月): 真实视频集成、数据持久化、高级分析",
            "中期目标(6-12个月): 社交功能、智能推荐、多语言支持",
            "长期目标(1-2年): VR/AR集成、AI教练、专业版功能",
            "具体规划: 集成真实世界杯进球视频、支持云端存储、多设备数据同步",
            "添加更多技能维度、对比专业运动员数据、生成专业训练报告",
            "训练社区、分享训练成果、挑战和竞赛、基于机器学习的训练推荐"
        ]
    )
    
    slide13 = add_content_slide(
        "总结",
        [
            "核心价值: 智能化(AI姿势识别)、数据化(量化训练数据)、个性化(定制训练计划)、游戏化(成就系统)",
            "技术优势: 现代化技术栈、AI驱动、优秀的UI/UX、响应式设计、数据安全",
            "应用场景: 个人运动员训练、教练指导学员、体育机构管理、青少年培训",
            "商业价值: 降低培训成本、提高训练效率、提升训练效果、增强用户体验",
            "让我们一起用AI技术，推动足球训练的数字化转型！"
        ]
    )
    
    slide14 = add_content_slide(
        "Q&A",
        [
            "Q1: 这个系统适合哪些用户？",
            "A: 适合所有足球运动员，包括职业运动员、业余爱好者、青少年球员。同时也适合足球教练和体育培训机构。",
            "Q2: 需要什么设备？",
            "A: 只需要一台电脑或手机，配备摄像头即可。系统支持Windows、Mac、iOS、Android等主流平台。",
            "Q3: 数据安全吗？",
            "A: 非常安全。所有数据存储在本地，用户完全控制数据共享。支持导出备份，数据不会上传到第三方服务器。"
        ]
    )
    
    slide15 = add_content_slide(
        "感谢",
        [
            "感谢大家抽出宝贵时间参加本次交流会。",
            "足球教练助手是一个持续迭代的项目，我们致力于通过AI技术帮助足球运动员提升训练效果。",
            "联系方式:",
            "邮箱: support@soccercoach.com",
            "GitHub: github.com/soccercoach",
            "网站: www.soccercoach.com",
            "微信: SoccerCoachHelper",
            "本项目完全开源，欢迎Star、Fork、贡献代码、报告Bug和问题。",
            "让我们一起用AI技术，推动足球训练的数字化转型！谢谢大家！"
        ]
    )
    
    return prs

if __name__ == "__main__":
    prs = create_presentation()
    prs.save('足球教练助手演示.pptx')
    print("PPT已生成: 足球教练助手演示.pptx")